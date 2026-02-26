# ==============================
# SMART CITY CHATBOT REPORT GENERATOR
# ==============================

import os
import pickle
from pptx import Presentation
from pptx.util import Inches
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.lib.pagesizes import A4

# ==============================
# Paths
# ==============================
output_dir = os.path.join(os.getcwd(), "report_files")
os.makedirs(output_dir, exist_ok=True)

ppt_path = os.path.join(output_dir, "Smart_City_Chatbot_Internship_Presentation.pptx")
pdf_path = os.path.join(output_dir, "Smart_City_Chatbot_Internship_Report.pdf")

# ==============================
# 1️⃣ CREATE POWERPOINT
# ==============================
prs = Presentation()

slides_content = [
    ("Smart City Chatbot System", "Internship Project Presentation\nGuwahati Smart City Assistant"),
    ("Introduction", "• ML-based Smart City Assistant\n• Provides city-related information\n• Built using Python & Streamlit"),
    ("Problem Statement", "• Citizens struggle to find city information\n• Manual search consumes time\n• Need for intelligent assistant"),
    ("Objectives", "• Develop ML intent classifier\n• Integrate ML with database\n• Build interactive web app"),
    ("Technologies Used", "• Python\n• Streamlit\n• Scikit-learn\n• SQLite\n• Pandas"),
    ("System Architecture", "User → Streamlit UI → ML Model → Database → Output"),
    ("Machine Learning Model", "• TF-IDF Vectorization\n• Logistic Regression\n• Intent Classification"),
    ("Database Design", "• Hospitals\n• Government Offices\n• Tourist Places\n• Utilities"),
    ("Implementation", "• Trained model\n• Saved model files\n• Integrated with Streamlit\n• Connected SQLite database"),
    ("Challenges Faced", "• Model path errors\n• Database issues\n• Intent mismatch\n• Debugging Streamlit errors"),
    ("Results", "• Correct intent detection\n• Accurate data retrieval\n• Fully working chatbot"),
    ("Learning Outcomes", "• ML model deployment\n• Database integration\n• Real-world debugging experience"),
    ("Conclusion", "Smart City Chatbot demonstrates integration of ML + Web Development for citizen service solutions"),
]

for title, content in slides_content:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

prs.save(ppt_path)
print(f"✅ PowerPoint created: {ppt_path}")

# ==============================
# 2️⃣ CREATE PDF
# ==============================
doc = SimpleDocTemplate(pdf_path, pagesize=A4)
elements = []
styles = getSampleStyleSheet()
title_style = styles["Heading1"]
normal_style = styles["Normal"]

elements.append(Paragraph("SMART CITY CHATBOT SYSTEM", title_style))
elements.append(Spacer(1, 0.3 * inch))

sections = [
    "Introduction: Smart City Chatbot provides city information using Machine Learning and Streamlit.",
    "Problem Statement: Citizens face difficulty accessing important city services quickly.",
    "Technologies Used: Python, Streamlit, Scikit-learn, SQLite, Pandas.",
    "System Architecture: User → UI → ML Model → Database → Output.",
    "Machine Learning: TF-IDF Vectorizer with Logistic Regression for intent classification.",
    "Database: SQLite database containing hospitals, offices, tourist places and utilities.",
    "Results: System successfully predicts intent and retrieves relevant information.",
    "Learning Outcomes: Hands-on experience in ML, database management and web app development.",
    "Conclusion: Project demonstrates real-world integration of AI and smart city services."
]

for section in sections:
    elements.append(Paragraph(section, normal_style))
    elements.append(Spacer(1, 0.3 * inch))

doc.build(elements)
print(f"✅ PDF created: {pdf_path}")
